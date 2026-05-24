"""
Fase 4d: Extrai texto e imagens de um PPTX, slide por slide.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def _processar_shape(shape, blocos, contador_imagens, pasta_saida, num_slide):
    """
    Processa um shape (forma) do slide. Pode ser texto, imagem, tabela ou grupo.
    Retorna o contador de imagens atualizado.
    """
    # 1) GRUPO: chama recursivamente pra cada shape interno
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            contador_imagens = _processar_shape(
                sub_shape, blocos, contador_imagens, pasta_saida, num_slide
            )
        return contador_imagens

    # 2) IMAGEM
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        try:
            imagem = shape.image
            ext = imagem.ext  # "png", "jpeg", etc.
            if ext == "jpeg":
                ext = "jpg"

            contador_imagens += 1
            nome_arquivo = f"img_{contador_imagens}.{ext}"
            caminho_imagem = os.path.join(pasta_saida, nome_arquivo)

            with open(caminho_imagem, "wb") as f:
                f.write(imagem.blob)

            blocos.append({
                "tipo": "imagem",
                "slide": num_slide,
                "caminho": caminho_imagem
            })
        except Exception as e:
            print(f"   ⚠️  Não consegui extrair imagem do slide {num_slide}: {e}")
        return contador_imagens

    # 3) TABELA
    if shape.has_table:
        linhas = []
        for linha in shape.table.rows:
            celulas = [celula.text.strip() for celula in linha.cells]
            if any(celulas):
                linhas.append(" | ".join(celulas))
        if linhas:
            texto_tabela = "Tabela: " + ". ".join(linhas)
            blocos.append({
                "tipo": "texto",
                "slide": num_slide,
                "conteudo": texto_tabela
            })
        return contador_imagens

    # 4) TEXTO
    if shape.has_text_frame:
        texto = shape.text_frame.text.strip()
        if texto:
            blocos.append({
                "tipo": "texto",
                "slide": num_slide,
                "conteudo": texto
            })
        return contador_imagens

    return contador_imagens


def extrair_conteudo_pptx(caminho_pptx: str, pasta_saida: str = "saida_imagens"):
    """
    Lê um PPTX e devolve uma lista de blocos na ordem dos slides.
    Inclui marcadores de slide pra orientar a pessoa cega.
    """
    Path(pasta_saida).mkdir(exist_ok=True)

    apresentacao = Presentation(caminho_pptx)
    blocos = []
    contador_imagens = 0
    total_slides = len(apresentacao.slides)

    for num_slide, slide in enumerate(apresentacao.slides, start=1):
        print(f"Processando slide {num_slide}/{total_slides}...")

        # Marcador de início do slide (vira texto que será narrado)
        blocos.append({
            "tipo": "texto",
            "slide": num_slide,
            "conteudo": f"Slide {num_slide} de {total_slides}."
        })

        # Ordena shapes por posição (top, depois left) pra leitura natural
        # Shapes sem posição definida vão pro fim
        def chave_ordem(s):
            try:
                return (s.top if s.top is not None else 999999,
                        s.left if s.left is not None else 999999)
            except Exception:
                return (999999, 999999)

        shapes_ordenados = sorted(slide.shapes, key=chave_ordem)

        for shape in shapes_ordenados:
            contador_imagens = _processar_shape(
                shape, blocos, contador_imagens, pasta_saida, num_slide
            )

        # Notas do apresentador (speaker notes) — muito importantes pra acessibilidade
        if slide.has_notes_slide:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                blocos.append({
                    "tipo": "texto",
                    "slide": num_slide,
                    "conteudo": f"Notas do apresentador: {notas}"
                })

    return blocos


# --- Teste rápido ---
if __name__ == "__main__":
    caminho = "teste.pptx"

    if not os.path.exists(caminho):
        print(f"❌ Arquivo '{caminho}' não encontrado.")
        print("   Coloque um PPTX chamado 'teste.pptx' na pasta do projeto.")
    else:
        blocos = extrair_conteudo_pptx(caminho)
        print(f"\n✅ Extração concluída! {len(blocos)} blocos encontrados.\n")

        for i, bloco in enumerate(blocos, start=1):
            if bloco["tipo"] == "texto":
                preview = bloco["conteudo"][:80].replace("\n", " ")
                print(f"[{i}] TEXTO  (slide {bloco['slide']}): {preview}...")
            else:
                print(f"[{i}] IMAGEM (slide {bloco['slide']}): {bloco['caminho']}")